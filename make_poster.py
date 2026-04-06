"""
4 letter-sized slides (8.5 x 11 inches, portrait) for VLM Visual Agnosia project.
One slide per topic: Motivation | Methods | Results | Conclusions
Run: python make_poster.py
Output: poster.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Letter size portrait
SW = Inches(8.5)
SH = Inches(11)

SAGE_DARK  = RGBColor(58,  88,  68)
SAGE_MID   = RGBColor(95,  133, 103)
SAGE_LIGHT = RGBColor(185, 213, 190)
SAGE_PALE  = RGBColor(238, 247, 240)
WHITE      = RGBColor(255, 255, 255)
BLACK      = RGBColor(18,  18,  18)
GRAY       = RGBColor(90,  90,  90)
DARK_TEXT  = RGBColor(28,  48,  33)
GOLD       = RGBColor(190, 158, 72)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Background
    s = slide.shapes.add_shape(1, 0, 0, SW, SH)
    s.fill.solid(); s.fill.fore_color.rgb = SAGE_PALE
    s.line.fill.background()
    return slide

def header(slide, title, subtitle=None):
    """Top header bar with title and optional subtitle."""
    hh = Inches(1.35) if subtitle else Inches(1.0)
    s = slide.shapes.add_shape(1, 0, 0, SW, hh)
    s.fill.solid(); s.fill.fore_color.rgb = SAGE_DARK
    s.line.fill.background()
    # gold accent
    g = slide.shapes.add_shape(1, 0, hh, SW, Inches(0.055))
    g.fill.solid(); g.fill.fore_color.rgb = GOLD
    g.line.fill.background()
    # title
    tb = slide.shapes.add_textbox(Inches(0.25), Inches(0.07), SW-Inches(0.5), Inches(0.62))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = WHITE
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.25), Inches(0.68), SW-Inches(0.5), Inches(0.5))
        tf2 = tb2.text_frame; tf2.word_wrap = True
        p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(13); r2.font.italic = True; r2.font.color.rgb = SAGE_LIGHT
    return hh + Inches(0.055)

def footer(slide, text):
    s = slide.shapes.add_shape(1, 0, SH-Inches(0.32), SW, Inches(0.32))
    s.fill.solid(); s.fill.fore_color.rgb = SAGE_DARK; s.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.2), SH-Inches(0.3), SW-Inches(0.4), Inches(0.28))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(9); r.font.color.rgb = SAGE_LIGHT

def box(slide, label, l, t, w, h):
    """Section box: sage header + white body."""
    bh = Inches(0.38)
    # border
    b = slide.shapes.add_shape(1, l, t, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = WHITE
    b.line.color.rgb = SAGE_DARK; b.line.width = Pt(1.5)
    # header
    hd = slide.shapes.add_shape(1, l, t, w, bh)
    hd.fill.solid(); hd.fill.fore_color.rgb = SAGE_MID; hd.line.fill.background()
    ac = slide.shapes.add_shape(1, l, t, Inches(0.055), bh)
    ac.fill.solid(); ac.fill.fore_color.rgb = GOLD; ac.line.fill.background()
    tb = slide.shapes.add_textbox(l+Inches(0.14), t+Inches(0.05), w-Inches(0.2), bh-Inches(0.06))
    tf = tb.text_frame; p = tf.paragraphs[0]
    r = p.add_run(); r.text = label
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
    return t + bh  # y of body start

def txt(slide, text, l, t, w, h, size=16, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def bullets(slide, items, l, t, w, h, size=16, color=BLACK, spacing=6):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.space_before = Pt(spacing)
        indent = item.startswith("  ")
        r = p.add_run()
        r.text = ("• " if not indent else "   ◦ ") + item.strip()
        r.font.size = Pt(size-1 if indent else size)
        r.font.color.rgb = GRAY if indent else color
    return tb

def image(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)

M  = Inches(0.3)   # margin
P  = Inches(0.15)  # inner padding
G  = Inches(0.14)  # gap between boxes
IW = SW - M*2      # inner width

FOOT = "Haya Bugshan  •  Mohammad Essa  •  CPSY 1950: Deep Learning in Brains, Minds & Machines  •  Brown University  •  Spring 2026"

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — MOTIVATION
# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH

sl1 = new_slide(prs)
y = header(sl1,
    "Do VLMs Have Human-Like Vision, or Just Pattern Matching?",
    "Motivation  •  Background  •  Hypothesis")
footer(sl1, FOOT)

BODY_BOT = SH - Inches(0.32) - G
y += G

# Box 1: Cognitive phenomenon
BH1 = Inches(2.9)
yb = box(sl1, "THE COGNITIVE PHENOMENON", M, y, IW, BH1)
bullets(sl1, [
    "Human object recognition involves two dissociable processing stages (Riddoch & Humphreys, 1987):",
    "  Apperceptive stage — the brain performs global form completion, assembling a coherent percept from raw visual input before any meaning is accessed",
    "  Associative stage — the completed percept is matched to stored semantic knowledge to yield recognition",
    "These stages can break independently (visual agnosia). Apperceptive patients cannot form percepts from degraded input; associative patients form percepts but cannot access meaning.",
    "Mooney (two-tone) images isolate the apperceptive stage — objects are invisible until the visual system commits to a global form hypothesis. Language cannot help before the percept forms.",
    "Degraded line drawings allow bottom-up feature matching without requiring global completion — the apperceptive stage is partially bypassed.",
], M+P, yb+P, IW-P*2, BH1-Inches(0.38)-P, size=16, spacing=6)

y += BH1 + G

# Box 2: Why test on VLMs
BH2 = Inches(2.0)
yb2 = box(sl1, "WHY TEST THIS ON VLMs?", M, y, IW, BH2)
bullets(sl1, [
    "Marjieh et al. (2023): adding visual input to GPT-4 does not improve perceptual similarity judgments — suggesting VLM vision may not add genuine perceptual grounding beyond language.",
    "If VLMs process vision through language-mediated feature matching (skipping the apperceptive stage), Mooney images should specifically expose that gap.",
    "This tests Marjieh's null result at the algorithmic level (Marr Level 2) — asking which processing stage breaks, not just measuring representational geometry.",
], M+P, yb2+P, IW-P*2, BH2-Inches(0.38)-P, size=16, spacing=6)

y += BH2 + G

# Box 3: Hypothesis (highlighted)
BH3 = BODY_BOT - y
yb3 = box(sl1, "HYPOTHESIS", M, y, IW, BH3)
s = sl1.shapes.add_shape(1, M+P, yb3+P, IW-P*2, BH3-Inches(0.38)-P*2)
s.fill.solid(); s.fill.fore_color.rgb = SAGE_LIGHT; s.line.fill.background()
txt(sl1,
    "VLMs will show higher accuracy on degraded line drawings than on Mooney images, "
    "reflecting a selective deficit at the apperceptive (global form completion) stage.\n\n"
    "If confirmed: VLM vision is texture/feature-based rather than genuinely perceptual — "
    "a behavioral-level mechanism for Marjieh et al.'s finding that visual input does not "
    "add perceptual structure beyond what language already encodes.",
    M+P+Inches(0.12), yb3+P+Inches(0.14),
    IW-P*2-Inches(0.2), BH3-Inches(0.38)-P*2-Inches(0.18),
    size=17, color=DARK_TEXT, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — METHODS
# ══════════════════════════════════════════════════════════════════════════════
sl2 = new_slide(prs)
y = header(sl2, "Methods", "Task design  •  Stimuli  •  Models  •  Prompting  •  Confound control")
footer(sl2, FOOT)
y += G

# Stimuli
BH1 = Inches(2.55)
yb = box(sl2, "STIMULI  (161 matched pairs)", M, y, IW, BH1)
bullets(sl2, [
    "Base set: Snodgrass & Vanderwart (1980) standardized line drawings — 161 common objects (Rossion & Pourtois, 2004 version)",
    "Degraded condition: 50% of dark line pixels randomly removed (NumPy random seed = 42). Partial features preserved; bottom-up matching possible.",
    "Mooney condition: Gaussian blur (σ = 3) applied, then binary threshold at pixel value 128 → pure black/white blobs. No legible edges or parts; global form completion required.",
    "Each object appears in both conditions — perfectly matched pairs eliminate object-level confounds.",
], M+P, yb+P, IW-P*2, BH1-Inches(0.38)-P, size=16, spacing=6)

y += BH1 + G

# Stimulus image strip
BH_IMG = Inches(2.1)
BASE = "/Users/hayabugshan/Desktop/vlm-visual-agnosia/stimuli"
ex_pairs = [
    ("goat", f"{BASE}/degraded/goat.png",    f"{BASE}/mooney/goat.png"),
    ("frog", f"{BASE}/degraded/frog.png",     f"{BASE}/mooney/frog.png"),
    ("giraffe", f"{BASE}/degraded/giraffe.png", f"{BASE}/mooney/giraffe.png"),
    ("glasses", f"{BASE}/degraded/glasses.png", f"{BASE}/mooney/glasses.png"),
]
yb_img = box(sl2, "STIMULUS EXAMPLES  (degraded left, mooney right per pair)", M, y, IW, BH_IMG)
n = len(ex_pairs)
cell_w = (IW - P*2) / (n * 2)
img_h = BH_IMG - Inches(0.38) - P - Inches(0.28)
for i, (label, dp, mp) in enumerate(ex_pairs):
    x = M + P + i * 2 * cell_w
    if os.path.exists(dp):
        sl2.shapes.add_picture(dp, x+Inches(0.02), yb_img+Inches(0.02), cell_w-Inches(0.04), img_h)
    if os.path.exists(mp):
        sl2.shapes.add_picture(mp, x+cell_w+Inches(0.02), yb_img+Inches(0.02), cell_w-Inches(0.04), img_h)
    txt(sl2, label, x, yb_img+img_h+Inches(0.04), cell_w*2, Inches(0.24),
        size=11, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

y += BH_IMG + G

# Models + Prompting
BH2 = Inches(1.9)
yb2 = box(sl2, "MODELS & PROMPTING", M, y, IW, BH2)
bullets(sl2, [
    "Models: gemini-3-flash-preview (open-weight, current results); claude-sonnet-4-5 & gemini-2.5-pro (frontier, pending after April 5)",
    "System prompt: \"Respond with only the name of the object you see — a single word or short noun phrase. Do not describe the image.\"",
    "User prompt: \"What object is in this image?\"  —  identical across all conditions and models (no chain-of-thought, reasoning_effort=None)",
    "Scoring: correct if true object label appears as a case-insensitive substring of the model response",
], M+P, yb2+P, IW-P*2, BH2-Inches(0.38)-P, size=16, spacing=6)

y += BH2 + G

# Confound control
BH3 = BODY_BOT - y
yb3 = box(sl2, "CONFOUND CONTROL", M, y, IW, BH3)
bullets(sl2, [
    "Same 161 objects in both conditions — eliminates object-level confounds between Mooney and degraded",
    "Identical prompt wording and structure across all models — eliminates autoregressive framing artifacts",
    "Fixed random seed (42) for pixel removal — degradation is reproducible across runs",
    "No chain-of-thought used — consistent prompting strategy prevents CoT-related distribution shift",
    "Limitation acknowledged: human baseline from Snodgrass & Corwin (1988) uses a different fragmentation paradigm — not directly comparable; treated as approximate",
], M+P, yb3+P, IW-P*2, BH3-Inches(0.38)-P, size=16, spacing=6)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — RESULTS
# ══════════════════════════════════════════════════════════════════════════════
sl3 = new_slide(prs)
y = header(sl3, "Results", "gemini-3-flash-preview  •  161 objects  •  Mooney vs. Degraded")
footer(sl3, FOOT)
y += G

# Main figure box
BH_FIG = Inches(5.8)
yb_fig = box(sl3, "ACCURACY BY CONDITION  —  VLM vs. Human Benchmark", M, y, IW, BH_FIG)
fig = "/Users/hayabugshan/Desktop/vlm-visual-agnosia/accuracy_with_human_baseline.png"
fig_h = BH_FIG - Inches(0.38) - Inches(0.72) - P
if os.path.exists(fig):
    sl3.shapes.add_picture(fig, M+P, yb_fig+Inches(0.08), IW-P*2, fig_h)
txt(sl3,
    "Figure 1. Proportion correct for gemini-3-flash-preview on Mooney (blue) vs. degraded (orange) conditions. "
    "Dashed lines = approximate human benchmarks from Snodgrass & Corwin (1988): ~34% Mooney, ~50% degraded. "
    "Human baselines reflect a different fragmentation paradigm — treat as approximate.",
    M+P, yb_fig+Inches(0.08)+fig_h+Inches(0.06), IW-P*2, Inches(0.65),
    size=12, color=GRAY, italic=True)

y += BH_FIG + G

# Key findings
BH2 = Inches(2.1)
yb2 = box(sl3, "KEY FINDINGS", M, y, IW, BH2)
bullets(sl3, [
    "gemini-3-flash-preview: 65% correct on Mooney images, 85% correct on degraded drawings",
    "Expected dissociation present: degraded > Mooney — same direction as humans",
    "20 pt VLM gap (degraded vs. Mooney) mirrors ~16 pt human dissociation — same processing bottleneck",
    "VLM Mooney accuracy exceeds human benchmark (~34%) — likely because line-drawing Mooneys are less ambiguous than photograph-based stimuli; see limitations",
], M+P, yb2+P, IW-P*2, BH2-Inches(0.38)-P, size=16, spacing=6)

y += BH2 + G

# Summary table
BH3 = BODY_BOT - y
yb3 = box(sl3, "SUMMARY TABLE", M, y, IW, BH3)
rows = [
    ("Condition",         "VLM Accuracy", "Human Baseline", "Direction"),
    ("Degraded drawings", "85%",          "~35–63%",        "VLM > Human"),
    ("Mooney images",     "65%",          "~34%",           "VLM > Human"),
    ("Dissociation Δ",    "20 pts",       "~16 pts",        "Same direction ✓"),
]
tw  = IW - P*2
th  = (BH3 - Inches(0.38) - P) / len(rows)
cws = [tw*0.34, tw*0.2, tw*0.22, tw*0.24]
tx  = M + P
for ri, row in enumerate(rows):
    cx = tx
    hdr = ri == 0
    bg  = SAGE_MID if hdr else (SAGE_LIGHT if ri%2==0 else WHITE)
    fc  = WHITE if hdr else BLACK
    for cell, cw in zip(row, cws):
        s = sl3.shapes.add_shape(1, cx, yb3+ri*th, cw, th)
        s.fill.solid(); s.fill.fore_color.rgb = bg
        s.line.color.rgb = SAGE_MID; s.line.width = Pt(0.5)
        tb = sl3.shapes.add_textbox(cx+Inches(0.07), yb3+ri*th+Inches(0.06),
                                    cw-Inches(0.1), th-Inches(0.08))
        tf = tb.text_frame; p = tf.paragraphs[0]
        r = p.add_run(); r.text = cell
        r.font.size = Pt(15); r.font.bold = hdr; r.font.color.rgb = fc
        cx += cw

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — OPEN QUESTIONS & CONCLUSIONS
# ══════════════════════════════════════════════════════════════════════════════
sl4 = new_slide(prs)
y = header(sl4, "Open Questions & Conclusions",
           "Interpretation  •  Marr's Levels  •  Limitations  •  Next Steps")
footer(sl4, FOOT)
y += G

# Interpretation
BH1 = Inches(2.2)
yb = box(sl4, "WHAT THE RESULTS TELL US", M, y, IW, BH1)
bullets(sl4, [
    "Dissociation (degraded > Mooney) is present in the predicted direction — consistent with VLMs finding global form completion harder than feature-based matching",
    "VLM vision appears predominantly texture/feature-based: succeeds when parts are physically present, struggles when global synthesis is required first",
    "Directionally consistent with Marjieh et al. — if VLMs skip the apperceptive stage, Mooney images should specifically expose that gap, and they do",
], M+P, yb+P, IW-P*2, BH1-Inches(0.38)-P, size=16, spacing=6)

y += BH1 + G

# Marr's levels
BH2 = Inches(1.95)
yb2 = box(sl4, "CONNECTION TO COURSE THEMES (MARR'S LEVELS)", M, y, IW, BH2)
lw3 = (IW - P*2) / 3
box_h = BH2 - Inches(0.38) - P*2
levels = [
    ("Level 1\nComputational", "What is vision for?\nObject recognition\nfrom images.", SAGE_DARK, WHITE),
    ("Level 2  ★\nAlgorithmic", "Which stage fails?\nApperceptive vs.\nassociative.\n← THIS PAPER", SAGE_MID, WHITE),
    ("Level 3\nImplementation", "Neural responses\nin IT / V4 cortex.\n(Brain-Score, Yamins)", SAGE_LIGHT, DARK_TEXT),
]
for i, (label, desc, bg, fc) in enumerate(levels):
    lx = M + P + i * lw3
    s = sl4.shapes.add_shape(1, lx+Inches(0.03), yb2+P, lw3-Inches(0.07), box_h)
    s.fill.solid(); s.fill.fore_color.rgb = bg
    s.line.color.rgb = SAGE_DARK; s.line.width = Pt(1)
    txt(sl4, label, lx+Inches(0.08), yb2+P+Inches(0.06),
        lw3-Inches(0.14), Inches(0.5), size=13, bold=True, color=fc, align=PP_ALIGN.CENTER)
    txt(sl4, desc, lx+Inches(0.08), yb2+P+Inches(0.56),
        lw3-Inches(0.14), box_h-Inches(0.6), size=13, color=fc, align=PP_ALIGN.CENTER)

y += BH2 + G

# Limitations
BH3 = Inches(1.95)
yb3 = box(sl4, "LIMITATIONS", M, y, IW, BH3)
bullets(sl4, [
    "Mooney stimuli from line drawings, not photographs — reduces blob ambiguity; likely underestimates task difficulty for outline-based objects (glasses, tools)",
    "Human baseline (Snodgrass & Corwin, 1988) uses a different fragmentation paradigm — approximate, not directly comparable",
    "Only one model tested so far; frontier model comparison (claude-sonnet-4-5, gemini-2.5-pro) pending",
    "Substring scoring may overcount correct responses for verbose or hedging outputs",
], M+P, yb3+P, IW-P*2, BH3-Inches(0.38)-P, size=16, spacing=6)

y += BH3 + G

# Future
BH4 = Inches(1.85)
yb4 = box(sl4, "NEXT STEPS", M, y, IW, BH4)
bullets(sl4, [
    "Add frontier models (claude-sonnet-4-5, gemini-2.5-pro) for model-to-model comparison",
    "Replace line-drawing Mooneys with MoonBase photograph-based stimuli (Tuerlinckx et al.) for a stronger manipulation",
    "Post-disambiguation condition: show model the original image first, then the Mooney — does recognition recover? (Mirrors Hegdé et al., 2007)",
    "Analyze token logprobs: is model uncertainty higher on Mooney than degraded?",
], M+P, yb4+P, IW-P*2, BH4-Inches(0.38)-P, size=16, spacing=6)

y += BH4 + G

# Conclusions
BH5 = BODY_BOT - y
yb5 = box(sl4, "CONCLUSIONS", M, y, IW, BH5)
s = sl4.shapes.add_shape(1, M+P, yb5+P, IW-P*2, BH5-Inches(0.38)-P*2)
s.fill.solid(); s.fill.fore_color.rgb = SAGE_LIGHT; s.line.fill.background()
txt(sl4,
    "gemini-3-flash-preview shows the predicted dissociation: degraded > Mooney. "
    "This is consistent with VLMs lacking genuine apperceptive completion — the pre-linguistic "
    "form assembly step that must precede semantic access in human vision.\n\n"
    "Our results provide a behavioral-level account of Marjieh et al.'s null result: "
    "if the apperceptive stage is absent, more visual input simply adds more of the same "
    "shallow features language already encodes.",
    M+P+Inches(0.12), yb5+P+Inches(0.14),
    IW-P*2-Inches(0.2), BH5-Inches(0.38)-P*2-Inches(0.18),
    size=17, color=DARK_TEXT, italic=True)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/hayabugshan/Desktop/vlm-visual-agnosia/poster.pptx"
prs.save(out)
print(f"Saved: {out}")
